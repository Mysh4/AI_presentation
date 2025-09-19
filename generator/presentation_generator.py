from pptx import Presentation
from pptx.util import Inches, Pt

class PresentationGenerator:
    def __init__(self, title, content_sections, conclusions, bibliography, num_main_slides=3):
        self.title = title
        self.content_sections = content_sections
        self.conclusions = conclusions
        self.bibliography = bibliography
        self.num_main_slides = num_main_slides
        self.prs = Presentation()

    def add_title_slide(self):
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = self.title
        slide.placeholders[1].text = "Автоматически сгенерировано ИИ"

    def add_content_slide(self):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Содержание"
        content = "\n".join([f"{i+1}. {section}" for i, section in enumerate(self.content_sections)])
        slide.placeholders[1].text = content

    def add_main_slides(self):
        slide_layout = self.prs.slide_layouts[1]
        for i in range(self.num_main_slides):
            slide = self.prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = self.content_sections[i] if i < len(self.content_sections) else f"Слайд {i+1}"
            slide.placeholders[1].text = f"Краткое описание: {self.content_sections[i]}"

    def add_conclusions_slide(self):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Выводы"
        slide.placeholders[1].text = self.conclusions

    def add_bibliography_slide(self):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Библиография"
        slide.placeholders[1].text = "\n".join(self.bibliography)

    def generate(self, filename="presentation.pptx"):
        self.add_title_slide()
        self.add_content_slide()
        self.add_main_slides()
        self.add_conclusions_slide()
        self.add_bibliography_slide()
        self.prs.save(filename)
        print(f"Презентация сохранена как {filename}")

if __name__ == "__main__":
    title = "Пример презентации"
    content_sections = ["Введение", "Методы", "Результаты"]
    conclusions = "Основные выводы по теме."
    bibliography = ["Книга 1", "Статья 2"]
    generator = PresentationGenerator(title, content_sections, conclusions, bibliography, num_main_slides=3)
    generator.generate()
